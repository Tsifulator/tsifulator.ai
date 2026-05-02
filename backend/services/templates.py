"""
tsifl Template Engine
Generates IB-grade Excel and PowerPoint files server-side.
No Office.js limitations. No locale bugs. Deterministic output every time.

Supported templates:
  - Trading Comps (Excel + PPT)
  - Precedent Transactions (Excel)  [coming]
  - DCF (Excel)                     [coming]
"""

import io
from typing import Optional

# ── Shared brand constants ─────────────────────────────────────────────────
NAVY      = "002366"   # IB dark navy — headers, title bars
NAVY_MED  = "1F3864"   # slightly lighter navy — alternating header
BLUE_LIGHT= "D6E4F0"   # light blue — median/mean row bg
BLUE_FILL = "EBF3FB"   # very light — subtle alternating rows
WHITE     = "FFFFFF"
DARK_TEXT = "1A1A2E"
GREY_LINE = "BDC3C7"
GREEN_CHK = "27AE60"
RED_NEG   = "E74C3C"
FONT_NAME = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
#  EXCEL — Trading Comps
# ══════════════════════════════════════════════════════════════════════════════

def generate_comp_table_xlsx(payload: dict) -> bytes:
    """
    Generate a fully-formatted IB trading comps table as .xlsx bytes.

    payload = {
      "title": str,
      "date":  str,          # "As of May 2026"
      "currency": str,       # "USD ($M)"
      "companies": [
        { "ticker", "name", "period",
          "revenue", "gross_margin", "ebitda", "ebitda_margin",
          "share_price", "market_cap", "ev",
          "pe", "ev_ebitda", "ev_revenue" }
      ]
    }
    """
    from openpyxl import Workbook
    from openpyxl.styles import (
        Font, PatternFill, Alignment, Border, Side, numbers
    )
    from openpyxl.utils import get_column_letter
    from statistics import median, mean

    wb = Workbook()
    ws = wb.active
    ws.title = "Trading Comps"

    companies = payload.get("companies", [])
    title     = payload.get("title", "Trading Comps")
    date_str  = payload.get("date", "")
    currency  = payload.get("currency", "USD ($M)")

    # ── Style helpers ─────────────────────────────────────────────────────
    def fill(hex_color):
        return PatternFill("solid", fgColor=hex_color)

    def font(bold=False, color=DARK_TEXT, size=10, name=FONT_NAME, italic=False):
        return Font(bold=bold, color=color, size=size, name=name, italic=italic)

    def align(h="left", v="center", wrap=False):
        return Alignment(horizontal=h, vertical=v, wrap_text=wrap)

    thin  = Side(style="thin",   color=GREY_LINE)
    thick = Side(style="medium", color=NAVY)

    def border(top=None, bottom=None, left=None, right=None):
        return Border(top=top, bottom=bottom, left=left, right=right)

    def all_border(style="thin"):
        s = Side(style=style, color=GREY_LINE)
        return Border(top=s, bottom=s, left=s, right=s)

    def header_border():
        return Border(
            top=Side(style="medium", color=NAVY),
            bottom=Side(style="medium", color=NAVY),
            left=Side(style="thin", color=WHITE),
            right=Side(style="thin", color=WHITE),
        )

    # ── Column layout ──────────────────────────────────────────────────────
    # A=Ticker, B=Company, C=Period, D=Revenue, E=Gross%, F=EBITDA,
    # G=EBITDA%, H=EV, I=Mkt Cap, J=EV/EBITDA, K=EV/Rev, L=P/E
    cols = ["A","B","C","D","E","F","G","H","I","J","K","L"]
    col_widths = [10, 26, 10, 12, 12, 12, 12, 12, 12, 12, 10, 10]
    for i, (col, w) in enumerate(zip(cols, col_widths)):
        ws.column_dimensions[col].width = w

    # Row heights
    ws.row_dimensions[1].height = 28
    ws.row_dimensions[2].height = 16
    ws.row_dimensions[4].height = 32

    # ── Row 1 — Title ──────────────────────────────────────────────────────
    ws.merge_cells("A1:L1")
    c = ws["A1"]
    c.value         = title
    c.font          = Font(bold=True, color=WHITE, size=14, name=FONT_NAME)
    c.fill          = fill(NAVY)
    c.alignment     = align("left", "center")
    c.border        = border(bottom=thick)

    # ── Row 2 — Subtitle ──────────────────────────────────────────────────
    ws.merge_cells("A2:C2")
    ws["A2"].value     = date_str
    ws["A2"].font      = font(italic=True, color="6B7280", size=9)
    ws["A2"].alignment = align("left", "center")

    ws.merge_cells("D2:L2")
    ws["D2"].value     = currency
    ws["D2"].font      = font(italic=True, color="6B7280", size=9)
    ws["D2"].alignment = align("right", "center")

    # ── Row 3 — blank ─────────────────────────────────────────────────────
    ws.row_dimensions[3].height = 6

    # ── Row 4 — Column headers ────────────────────────────────────────────
    headers = [
        "Ticker", "Company", "Period",
        "Revenue", "Gross Margin", "EBITDA", "EBITDA Margin",
        "EV ($B)", "Mkt Cap ($B)",
        "EV / EBITDA", "EV / Revenue", "P / E"
    ]
    for col_letter, header_text in zip(cols, headers):
        c = ws[f"{col_letter}4"]
        c.value     = header_text
        c.font      = Font(bold=True, color=WHITE, size=9, name=FONT_NAME)
        c.fill      = fill(NAVY)
        c.alignment = align("center", "center", wrap=True)
        c.border    = header_border()

    # ── Data rows ──────────────────────────────────────────────────────────
    FMT_INT    = '#,##0'
    FMT_DEC1   = '#,##0.0'
    FMT_MULT   = '#,##0.0x'
    FMT_PCT    = '0.0%'
    FMT_PRICE  = '[$$-409]#,##0.00'

    ev_ebitda_vals, ev_rev_vals, pe_vals = [], [], []
    ebitda_margin_vals, gross_margin_vals = [], []

    for i, co in enumerate(companies):
        row = 5 + i
        ws.row_dimensions[row].height = 18
        row_fill = fill(BLUE_FILL) if i % 2 == 1 else fill(WHITE)

        def cell(col_letter, value, fmt=None, h_align="right"):
            c = ws[f"{col_letter}{row}"]
            c.value     = value
            c.font      = font(size=9)
            c.fill      = row_fill
            c.alignment = align(h_align, "center")
            c.border    = all_border()
            if fmt:
                c.number_format = fmt

        # Text cols
        cell("A", co.get("ticker",""), h_align="center")
        cell("B", co.get("name",""),   h_align="left")
        cell("C", co.get("period",""), h_align="center")
        ws[f"A{row}"].font = Font(bold=True, size=9, name=FONT_NAME)

        # Numeric cols
        cell("D", co.get("revenue"),       FMT_DEC1)
        cell("E", co.get("gross_margin"),  FMT_PCT)
        cell("F", co.get("ebitda"),        FMT_DEC1)
        cell("G", co.get("ebitda_margin"), FMT_PCT)
        cell("H", co.get("ev"),            FMT_DEC1)
        cell("I", co.get("market_cap"),    FMT_DEC1)
        cell("J", co.get("ev_ebitda"),     FMT_MULT)
        cell("K", co.get("ev_revenue"),    FMT_MULT)
        cell("L", co.get("pe"),            FMT_MULT)

        # Collect for stats
        if co.get("ev_ebitda"):  ev_ebitda_vals.append(co["ev_ebitda"])
        if co.get("ev_revenue"): ev_rev_vals.append(co["ev_revenue"])
        if co.get("pe"):         pe_vals.append(co["pe"])
        if co.get("ebitda_margin"): ebitda_margin_vals.append(co["ebitda_margin"])
        if co.get("gross_margin"):  gross_margin_vals.append(co["gross_margin"])

    # ── Median / Mean rows ─────────────────────────────────────────────────
    stat_rows = [
        ("Median", lambda vals: median(vals) if vals else None),
        ("Mean",   lambda vals: mean(vals)   if vals else None),
    ]
    blank_row  = 5 + len(companies)
    ws.row_dimensions[blank_row].height = 8

    for j, (label, fn) in enumerate(stat_rows):
        row = 5 + len(companies) + 1 + j
        ws.row_dimensions[row].height = 18

        def stat_cell(col_letter, value, fmt=None, bold=False, h="right"):
            c = ws[f"{col_letter}{row}"]
            c.value     = value
            c.font      = Font(bold=bold, color=DARK_TEXT, size=9, name=FONT_NAME)
            c.fill      = fill(BLUE_LIGHT)
            c.alignment = align(h, "center")
            c.border    = Border(
                top=Side(style="thin", color=NAVY),
                bottom=Side(style="thin", color=NAVY),
                left=Side(style="thin", color=GREY_LINE),
                right=Side(style="thin", color=GREY_LINE),
            )
            if fmt: c.number_format = fmt

        stat_cell("A", label,                bold=True, h="center")
        stat_cell("B", "",                   h="left")
        stat_cell("C", "")
        stat_cell("D", fn([c.get("revenue") for c in companies if c.get("revenue")]),      FMT_DEC1)
        stat_cell("E", fn(gross_margin_vals),  FMT_PCT)
        stat_cell("F", fn([c.get("ebitda") for c in companies if c.get("ebitda")]),        FMT_DEC1)
        stat_cell("G", fn(ebitda_margin_vals), FMT_PCT)
        stat_cell("H", fn([c.get("ev") for c in companies if c.get("ev")]),                FMT_DEC1)
        stat_cell("I", fn([c.get("market_cap") for c in companies if c.get("market_cap")]),FMT_DEC1)
        stat_cell("J", fn(ev_ebitda_vals),    FMT_MULT, bold=True)
        stat_cell("K", fn(ev_rev_vals),       FMT_MULT, bold=True)
        stat_cell("L", fn(pe_vals),           FMT_MULT, bold=True)

    # ── Sources & Methodology ──────────────────────────────────────────────
    src_row = 5 + len(companies) + 4
    ws.row_dimensions[src_row].height = 14
    ws.merge_cells(f"A{src_row}:L{src_row}")
    ws[f"A{src_row}"].value     = "Sources & Methodology"
    ws[f"A{src_row}"].font      = Font(bold=True, color=DARK_TEXT, size=8, name=FONT_NAME)
    ws[f"A{src_row}"].alignment = align("left", "center")

    for k, co in enumerate(companies):
        r = src_row + 1 + k
        ws.row_dimensions[r].height = 13
        ws.merge_cells(f"A{r}:C{r}")
        ws[f"A{r}"].value     = co.get("ticker","")
        ws[f"A{r}"].font      = Font(bold=True, size=8, name=FONT_NAME)
        ws[f"A{r}"].alignment = align("left", "center")
        ws.merge_cells(f"D{r}:H{r}")
        ws[f"D{r}"].value     = "Financial Modeling Prep (FMP)"
        ws[f"D{r}"].font      = Font(size=8, name=FONT_NAME, color="6B7280")
        ws[f"I{r}"].value     = co.get("period","")
        ws[f"I{r}"].font      = Font(size=8, name=FONT_NAME, color="6B7280")

    # ── Freeze panes at row 5 ─────────────────────────────────────────────
    ws.freeze_panes = "D5"

    # ── Print settings ─────────────────────────────────────────────────────
    ws.page_setup.orientation  = "landscape"
    ws.page_setup.fitToPage    = True
    ws.page_setup.fitToWidth   = 1
    ws.page_setup.fitToHeight  = 0

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT — Trading Comps Slide
# ══════════════════════════════════════════════════════════════════════════════

def generate_comp_slide_pptx(payload: dict) -> bytes:
    """
    Generate an IB-grade trading comps PPT slide as .pptx bytes.
    Returns a full presentation (1 comp slide) — append more slides as needed.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt
    import copy
    from lxml import etree
    from statistics import median, mean

    def rgb(hex_str):
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    companies = payload.get("companies", [])
    title     = payload.get("title", "Trading Comps")
    date_str  = payload.get("date", "")
    currency  = payload.get("currency", "USD")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)

    # ── Title bar ─────────────────────────────────────────────────────────
    from pptx.util import Inches, Pt, Emu
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(12.73), Inches(0.55)
    )
    tf = title_box.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold   = True
    run.font.size   = Pt(16)
    run.font.color.rgb = rgb(NAVY)
    run.font.name   = "Calibri"

    # Date subtitle
    date_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.72), Inches(6), Inches(0.3)
    )
    tf2 = date_box.text_frame
    p2  = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = f"{date_str}  |  {currency}"
    run2.font.size  = Pt(8)
    run2.font.color.rgb = rgb("6B7280")
    run2.font.name  = "Calibri"
    run2.font.italic = True

    # Thin navy divider line under title
    from pptx.util import Emu
    from pptx.oxml import parse_xml
    connector = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(0.3), Inches(0.95),
        Inches(13.03), Inches(0.95)
    )
    connector.line.color.rgb = rgb(NAVY)
    connector.line.width = Pt(1.5)

    # ── Build table ────────────────────────────────────────────────────────
    col_headers = [
        "Ticker", "Company", "Period",
        "Revenue", "Gross\nMargin", "EBITDA", "EBITDA\nMargin",
        "EV ($B)", "Mkt Cap\n($B)", "EV /\nEBITDA", "EV / Rev", "P / E"
    ]
    col_widths_in = [0.75, 2.2, 0.75, 0.85, 0.85, 0.85, 0.85, 0.85, 0.85, 0.85, 0.75, 0.7]

    n_data   = len(companies)
    n_rows   = 1 + n_data + 2 + 1  # header + data + median/mean + spacer
    n_cols   = len(col_headers)

    tbl_top    = Inches(1.05)
    tbl_left   = Inches(0.3)
    tbl_width  = Inches(12.73)
    tbl_height = Inches(0.32 * n_rows)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, tbl_left, tbl_top, tbl_width, tbl_height
    )
    tbl = tbl_shape.table

    # Column widths
    total_w_in = sum(col_widths_in)
    for ci, w in enumerate(col_widths_in):
        tbl.columns[ci].width = int(Inches(w * 12.73 / total_w_in))

    def set_cell(ri, ci, text, bold=False, font_size=8,
                 bg=None, fg=DARK_TEXT, align=PP_ALIGN.RIGHT,
                 italic=False):
        cell = tbl.cell(ri, ci)
        cell.text = ""
        tf = cell.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text) if text is not None else "—"
        run.font.bold   = bold
        run.font.size   = Pt(font_size)
        run.font.color.rgb = rgb(fg)
        run.font.name   = "Calibri"
        run.font.italic = italic
        if bg:
            fill_xml = f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:srgbClr val="{bg}"/></a:solidFill>'
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            # remove existing fill
            for old in tcPr.findall(qn('a:solidFill')):
                tcPr.remove(old)
            for old in tcPr.findall(qn('a:noFill')):
                tcPr.remove(old)
            tcPr.insert(0, parse_xml(fill_xml))

    def fmt_mult(v):
        if v is None: return "—"
        try: return f"{float(v):.1f}x"
        except: return "—"

    def fmt_dec(v, dp=1):
        if v is None: return "—"
        try: return f"{float(v):,.{dp}f}"
        except: return "—"

    def fmt_pct(v):
        if v is None: return "—"
        try: return f"{float(v)*100:.1f}%"
        except: return "—"

    # Header row
    for ci, h in enumerate(col_headers):
        halign = PP_ALIGN.CENTER if ci < 3 else PP_ALIGN.RIGHT
        set_cell(0, ci, h, bold=True, font_size=8,
                 bg=NAVY, fg=WHITE, align=halign)

    # Data rows
    for i, co in enumerate(companies):
        ri  = i + 1
        bg  = BLUE_FILL if i % 2 == 1 else WHITE
        set_cell(ri, 0, co.get("ticker",""),          bold=True, bg=bg, align=PP_ALIGN.CENTER)
        set_cell(ri, 1, co.get("name",""),            bg=bg, align=PP_ALIGN.LEFT)
        set_cell(ri, 2, co.get("period",""),          bg=bg, align=PP_ALIGN.CENTER, italic=True)
        set_cell(ri, 3, fmt_dec(co.get("revenue")),   bg=bg)
        set_cell(ri, 4, fmt_pct(co.get("gross_margin")), bg=bg)
        set_cell(ri, 5, fmt_dec(co.get("ebitda")),    bg=bg)
        set_cell(ri, 6, fmt_pct(co.get("ebitda_margin")), bg=bg)
        set_cell(ri, 7, fmt_dec(co.get("ev")),        bg=bg)
        set_cell(ri, 8, fmt_dec(co.get("market_cap")),bg=bg)
        set_cell(ri, 9, fmt_mult(co.get("ev_ebitda")),bg=bg, bold=True)
        set_cell(ri,10, fmt_mult(co.get("ev_revenue")),bg=bg)
        set_cell(ri,11, fmt_mult(co.get("pe")),       bg=bg)

    # Spacer row
    spacer_ri = 1 + n_data
    for ci in range(n_cols):
        set_cell(spacer_ri, ci, "", bg=WHITE)
    tbl.rows[spacer_ri].height = int(Inches(0.08))

    # Median / Mean
    def stat_vals(key):
        return [co[key] for co in companies if co.get(key) is not None]

    stat_defs = [
        ("Median", lambda k: median(stat_vals(k)) if stat_vals(k) else None),
        ("Mean",   lambda k: mean(stat_vals(k))   if stat_vals(k) else None),
    ]
    for j, (label, fn) in enumerate(stat_defs):
        ri = 1 + n_data + 1 + j
        set_cell(ri, 0, label,        bold=True, bg=BLUE_LIGHT, align=PP_ALIGN.CENTER)
        set_cell(ri, 1, "",           bg=BLUE_LIGHT, align=PP_ALIGN.LEFT)
        set_cell(ri, 2, "",           bg=BLUE_LIGHT)
        set_cell(ri, 3, fmt_dec(fn("revenue")),       bg=BLUE_LIGHT)
        set_cell(ri, 4, fmt_pct(fn("gross_margin")),  bg=BLUE_LIGHT)
        set_cell(ri, 5, fmt_dec(fn("ebitda")),        bg=BLUE_LIGHT)
        set_cell(ri, 6, fmt_pct(fn("ebitda_margin")), bg=BLUE_LIGHT)
        set_cell(ri, 7, fmt_dec(fn("ev")),            bg=BLUE_LIGHT)
        set_cell(ri, 8, fmt_dec(fn("market_cap")),    bg=BLUE_LIGHT)
        set_cell(ri, 9, fmt_mult(fn("ev_ebitda")),    bg=BLUE_LIGHT, bold=True)
        set_cell(ri,10, fmt_mult(fn("ev_revenue")),   bg=BLUE_LIGHT)
        set_cell(ri,11, fmt_mult(fn("pe")),           bg=BLUE_LIGHT)

    # ── Footnote ──────────────────────────────────────────────────────────
    fn_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(7.1), Inches(12.73), Inches(0.3)
    )
    tf3 = fn_box.text_frame
    p3  = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = f"Source: Financial Modeling Prep (FMP), Polygon.io  |  Generated by tsifl  |  {date_str}"
    run3.font.size   = Pt(7)
    run3.font.color.rgb = rgb("9CA3AF")
    run3.font.name   = "Calibri"
    run3.font.italic = True

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
#  Helper — build payload from live market data
# ══════════════════════════════════════════════════════════════════════════════

async def build_comp_payload(tickers: list[str], title: str = None) -> dict:
    """
    Fetch live data for a list of tickers and return a payload
    ready for generate_comp_table_xlsx / generate_comp_slide_pptx.

    Data flow:
      1. Polygon → price, market_cap_B, shares_outstanding
      2. FMP → revenue_ltm_M, ebitda_ltm_M, margins, net_debt_M
         (yfinance fallback if FMP rate-limited or no key)
      3. Compute EV = market_cap + net_debt
         EV/EBITDA, EV/Revenue, P/E from combined data
    """
    import asyncio
    from datetime import datetime
    from services.fmp import get_fundamentals as fmp_get
    from services.yfinance_service import get_fundamentals as yf_get
    from services.polygon import get_stocks_batch

    upper = [t.upper() for t in tickers]

    # ── 1. Fetch price + market cap from Polygon (fast, reliable) ──────────
    mkt_list = await get_stocks_batch(upper)
    mkt = {m["ticker"]: m for m in mkt_list if not m.get("error")}

    # ── 2. Fetch fundamentals (yfinance primary — free, no rate limit)
    #       FMP only if yfinance fails AND FMP key is set
    async def get_fund(ticker):
        d = await yf_get(ticker)
        if d.get("error"):
            d = await fmp_get(ticker)
        return ticker, d

    fund_results = await asyncio.gather(*[get_fund(t) for t in upper])
    fund = {t: d for t, d in fund_results}

    # ── 3. Assemble payload ──────────────────────────────────────────────────
    companies = []
    for ticker in upper:
        m = mkt.get(ticker, {})
        f = fund.get(ticker, {})

        if f.get("error") and not m:
            continue  # no data at all — skip

        # Price + market cap from Polygon
        price     = _safe_float(m.get("price"))
        mkt_cap_b = _safe_float(m.get("market_cap_B"))

        # Fundamentals — both FMP and yfinance use _ltm_M / _pct naming
        rev_m    = _safe_float(f.get("revenue_ltm_M"))
        ebitda_m = _safe_float(f.get("ebitda_ltm_M"))
        ni_m     = _safe_float(f.get("net_income_ltm_M"))
        gm_pct   = _safe_float(f.get("gross_margin_pct"))    # 0-100
        ebitda_margin_pct = (
            round(ebitda_m / rev_m * 100, 1)
            if ebitda_m and rev_m else None
        )
        net_debt_m = _safe_float(f.get("net_debt_M"))

        # EV = market_cap_B + net_debt_B
        ev_b = None
        if mkt_cap_b is not None and net_debt_m is not None:
            ev_b = round(mkt_cap_b + net_debt_m / 1e3, 2)
        elif mkt_cap_b is not None:
            ev_b = mkt_cap_b  # approximate if no debt data

        # Multiples
        ev_ebitda = round(ev_b * 1e3 / ebitda_m, 1) if ev_b and ebitda_m and ebitda_m > 0 else None
        ev_rev    = round(ev_b * 1e3 / rev_m,    1) if ev_b and rev_m    and rev_m    > 0 else None
        pe        = round(mkt_cap_b * 1e3 / ni_m, 1) if mkt_cap_b and ni_m and ni_m > 0 else None

        # Period label
        period = f.get("period_label") or "LTM"

        companies.append({
            "ticker":        ticker,
            "name":          f.get("name") or m.get("name") or ticker,
            "period":        period,
            "revenue":       round(rev_m / 1e3, 2) if rev_m else None,   # → $B for display
            "gross_margin":  gm_pct / 100 if gm_pct else None,           # → decimal
            "ebitda":        round(ebitda_m / 1e3, 2) if ebitda_m else None,
            "ebitda_margin": ebitda_margin_pct / 100 if ebitda_margin_pct else None,
            "share_price":   price,
            "market_cap":    mkt_cap_b,
            "ev":            ev_b,
            "pe":            pe,
            "ev_ebitda":     ev_ebitda,
            "ev_revenue":    ev_rev,
        })

    return {
        "title":     title or f"Trading Comps — {', '.join(upper)}",
        "date":      f"As of {datetime.utcnow().strftime('%B %Y')}",
        "currency":  "USD ($B)",
        "companies": companies,
    }


# ══════════════════════════════════════════════════════════════════════════════
#  POWERPOINT — Full 4-slide comp deck
# ══════════════════════════════════════════════════════════════════════════════

def generate_comp_deck_pptx(payload: dict) -> bytes:
    """
    Generate a full 4-slide IB comp deck:
      Slide 1 — Title
      Slide 2 — Trading Comps table
      Slide 3 — Multiples snapshot (big KPI numbers)
      Slide 4 — Key takeaways + sources
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml import parse_xml
    from pptx.oxml.ns import qn
    from statistics import median, mean

    def rgb(hex_str):
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    companies = payload.get("companies", [])
    title     = payload.get("title", "Trading Comps")
    date_str  = payload.get("date", "")
    currency  = payload.get("currency", "USD")

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # ── Force white background on all slides ──────────────────────────────
    def _set_white_bg(slide):
        from pptx.oxml.ns import qn as _qn
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # ── Shared helpers ────────────────────────────────────────────────────
    def add_textbox(slide, text, left, top, width, height,
                    bold=False, size=11, color=DARK_TEXT, align=PP_ALIGN.LEFT,
                    italic=False, bg=None):
        box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold    = bold
        run.font.size    = Pt(size)
        run.font.color.rgb = rgb(color)
        run.font.name    = "Calibri"
        run.font.italic  = italic
        if bg:
            fill_xml = (
                f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:srgbClr val="{bg}"/></a:solidFill>'
            )
            sp = box._element
            spPr = sp.find(qn('p:spPr'))
            if spPr is not None:
                for old in spPr.findall(qn('a:solidFill')):
                    spPr.remove(old)
                spPr.insert(0, parse_xml(fill_xml))
        return box

    def add_divider(slide, y=0.9):
        """Thin navy horizontal rule."""
        conn = slide.shapes.add_connector(
            1, Inches(0.3), Inches(y), Inches(13.03), Inches(y)
        )
        conn.line.color.rgb = rgb(NAVY)
        conn.line.width = Pt(1.2)

    def add_slide_header(slide, title_text, subtitle_text=""):
        add_textbox(slide, title_text,
                    0.3, 0.18, 12.73, 0.55,
                    bold=True, size=15, color=NAVY)
        if subtitle_text:
            add_textbox(slide, subtitle_text,
                        0.3, 0.7, 12.73, 0.25,
                        size=8, color="6B7280", italic=True)
        add_divider(slide, y=0.92)

    def add_footnote(slide, text=None):
        note = text or f"Source: Financial Modeling Prep (FMP), Polygon.io  |  Generated by tsifl  |  {date_str}"
        add_textbox(slide, note,
                    0.3, 7.1, 12.73, 0.3,
                    size=7, color="9CA3AF", italic=True)

    def fmt_mult(v):
        try: return f"{float(v):.1f}x" if v is not None else "—"
        except: return "—"

    def fmt_dec(v, dp=1):
        try: return f"{float(v):,.{dp}f}" if v is not None else "—"
        except: return "—"

    def fmt_pct(v):
        try: return f"{float(v)*100:.1f}%" if v is not None else "—"
        except: return "—"

    def stat_vals(key):
        return [co[key] for co in companies if co.get(key) is not None]

    def med(key):
        vals = stat_vals(key)
        return median(vals) if vals else None

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 1 — Title
    # ──────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank); _set_white_bg(s1)

    # Full navy top bar
    bar = s1.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0), Inches(13.33), Inches(2.2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(NAVY)
    bar.line.fill.background()

    add_textbox(s1, title, 0.5, 0.45, 12.33, 0.9,
                bold=True, size=32, color=WHITE)
    add_textbox(s1, f"Trading Comparables Analysis  ·  {date_str}",
                0.5, 1.35, 12.33, 0.4,
                size=13, color="A8C4E0", italic=True)

    # Company ticker chips
    if companies:
        chip_x = 0.5
        for co in companies[:8]:
            ticker = co.get("ticker","")
            chip = s1.shapes.add_shape(
                1, Inches(chip_x), Inches(2.7), Inches(0.9), Inches(0.35)
            )
            chip.fill.solid()
            chip.fill.fore_color.rgb = rgb("1F3864")
            chip.line.color.rgb = rgb("2E5599")
            chip.line.width = Pt(0.5)
            tf = chip.text_frame
            tf.text = ticker
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].runs[0]
            run.font.bold  = True
            run.font.size  = Pt(9)
            run.font.color.rgb = rgb(WHITE)
            run.font.name  = "Calibri"
            chip_x += 1.05

    add_textbox(s1,
                "Confidential  ·  For Discussion Purposes Only",
                0.5, 6.9, 12.33, 0.35,
                size=8, color="6B7280", italic=True,
                align=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 2 — Comp table (reuse single-slide generator)
    # ──────────────────────────────────────────────────────────────────────
    prs2_bytes = generate_comp_slide_pptx(payload)
    prs2 = Presentation(io.BytesIO(prs2_bytes))
    # Copy the slide XML directly
    template_slide = prs2.slides[0]
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
    import copy

    def _clone_slide(src_prs, src_slide, dest_prs):
        """Clone a slide from src_prs into dest_prs."""
        blank_layout = dest_prs.slide_layouts[6]
        new_slide = dest_prs.slides.add_slide(blank_layout)
        # Copy all shapes
        for shape in src_slide.shapes:
            el = copy.deepcopy(shape._element)
            new_slide.shapes._spTree.insert(2, el)
        # Copy background
        bg = src_slide.background
        new_bg = new_slide.background
        fill = bg.fill
        if fill.type is not None:
            new_bg.fill.solid()
            try:
                new_bg.fill.fore_color.rgb = fill.fore_color.rgb
            except Exception:
                pass
        return new_slide

    _clone_slide(prs2, template_slide, prs)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 3 — Multiples snapshot (big KPI boxes)
    # ──────────────────────────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank); _set_white_bg(s3)
    add_slide_header(s3, "Valuation Snapshot",
                     f"Median multiples — {date_str}  ·  {currency}")

    kpis = [
        ("EV / EBITDA",   fmt_mult(med("ev_ebitda")),  NAVY),
        ("EV / Revenue",  fmt_mult(med("ev_revenue")), "1F3864"),
        ("P / E",         fmt_mult(med("pe")),         "2E5599"),
    ]
    box_w, box_h = 3.6, 2.4
    spacing      = 0.4
    total_w      = len(kpis) * box_w + (len(kpis)-1) * spacing
    start_x      = (13.33 - total_w) / 2

    for i, (label, value, color) in enumerate(kpis):
        bx = start_x + i * (box_w + spacing)
        by = 1.8

        bg_shape = s3.shapes.add_shape(
            1, Inches(bx), Inches(by), Inches(box_w), Inches(box_h)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = rgb(color)
        bg_shape.line.fill.background()

        # Big number
        add_textbox(s3, value,
                    bx + 0.15, by + 0.25, box_w - 0.3, 1.3,
                    bold=True, size=44, color=WHITE,
                    align=PP_ALIGN.CENTER)
        # Label
        add_textbox(s3, label,
                    bx + 0.15, by + 1.65, box_w - 0.3, 0.45,
                    size=11, color="A8C4E0",
                    align=PP_ALIGN.CENTER, italic=True)

    # Company count note
    n = len(companies)
    add_textbox(s3,
                f"Based on {n} comparable compan{'y' if n==1 else 'ies'}  ·  Median shown",
                0.3, 4.55, 12.73, 0.3,
                size=8, color="9CA3AF", italic=True, align=PP_ALIGN.CENTER)

    add_footnote(s3)

    # ──────────────────────────────────────────────────────────────────────
    # SLIDE 4 — Key takeaways
    # ──────────────────────────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank); _set_white_bg(s4)
    add_slide_header(s4, "Key Observations", date_str)

    # Auto-generate observations from the data
    observations = []
    ev_vals  = stat_vals("ev_ebitda")
    rev_vals = stat_vals("ev_revenue")
    pe_vals  = stat_vals("pe")

    if ev_vals:
        hi = max(companies, key=lambda c: c.get("ev_ebitda") or 0)
        lo = min(companies, key=lambda c: c.get("ev_ebitda") or 999)
        observations.append(
            f"EV/EBITDA range: {fmt_mult(lo.get('ev_ebitda'))} ({lo.get('ticker','')}) "
            f"to {fmt_mult(hi.get('ev_ebitda'))} ({hi.get('ticker','')}), "
            f"median {fmt_mult(med('ev_ebitda'))}"
        )
    if pe_vals:
        observations.append(
            f"P/E multiples range {fmt_mult(min(pe_vals))}–{fmt_mult(max(pe_vals))}, "
            f"median {fmt_mult(med('pe'))}"
        )
    margins = stat_vals("ebitda_margin")
    if margins:
        best = max(companies, key=lambda c: c.get("ebitda_margin") or 0)
        observations.append(
            f"{best.get('name', best.get('ticker',''))} leads on EBITDA margin at "
            f"{fmt_pct(best.get('ebitda_margin'))}"
        )
    if len(companies) >= 2:
        by_ev = sorted([c for c in companies if c.get("ev_ebitda")],
                       key=lambda c: c["ev_ebitda"])
        if len(by_ev) >= 2:
            observations.append(
                f"{by_ev[0].get('name', by_ev[0].get('ticker',''))} trades at a discount "
                f"({fmt_mult(by_ev[0].get('ev_ebitda'))} EV/EBITDA) vs. "
                f"{by_ev[-1].get('name', by_ev[-1].get('ticker',''))} "
                f"({fmt_mult(by_ev[-1].get('ev_ebitda'))} EV/EBITDA)"
            )
    if not observations:
        observations = [
            "See trading comps table on prior slide for full detail",
            f"Analysis covers {len(companies)} publicly traded comparable companies",
            "All figures sourced from Financial Modeling Prep (FMP)"
        ]

    bullet_y = 1.15
    for obs in observations[:5]:
        # Bullet marker
        add_textbox(s4, "▸", 0.4, bullet_y, 0.3, 0.45,
                    bold=True, size=10, color=NAVY)
        add_textbox(s4, obs, 0.75, bullet_y, 11.8, 0.45,
                    size=10.5, color=DARK_TEXT)
        bullet_y += 0.7

    add_footnote(s4)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _safe_float(v, default=None):
    try: return round(float(v), 2) if v is not None else default
    except: return default

def _safe_pct(v, default=None):
    """Return as decimal (0.35 not 35) — formatters handle display."""
    try:
        f = float(v)
        return f / 100 if f > 1 else f  # normalise if stored as 35.0
    except: return default

def _safe_billions(v, default=None):
    try:
        f = float(v)
        return round(f / 1e9, 2) if f > 1e6 else round(f, 2)
    except: return default
